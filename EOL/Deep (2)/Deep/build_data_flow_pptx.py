"""
build_data_flow_pptx.py
=======================
Generate an editable PowerPoint flowchart in the classic process-
flow-chart style (rounded rectangles, ovals, arrows, light pastel
fills) showing the MES data flow:

   EOL Machine (oval)
        ↓
   PLC / Camera / Andon Box (rectangles)
        ↓
   Collector (rectangle)
        ↓
   5 data streams + Side Panel
        ↓
   Departments (coloured ovals — only colour layer)

All shapes are native PowerPoint shapes — click any shape, edit text,
move it around, recolour.  Save: D:\EOL\EOL\Deep (2)\Deep\MES_Data_Flow.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text  import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color  import RGBColor

# ─────────────────────────────────────────────────────────────────────
# Slide setup — 16:9 landscape, large enough for an A3 print
# ─────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(16.5)   # A3 landscape
prs.slide_height = Inches(11.7)

slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank layout

# Convenience for placing things in inches
def Pos(left, top, width, height):
    return (Inches(left), Inches(top), Inches(width), Inches(height))

# ─────────────────────────────────────────────────────────────────────
# Colour palette
# ─────────────────────────────────────────────────────────────────────
GREEN_LIGHT = RGBColor(0x52, 0xC4, 0x1A)   # green oval (start / source)
BLUE_DARK   = RGBColor(0x1F, 0x49, 0x7D)   # blue rect (key process)
ORANGE      = RGBColor(0xF5, 0x9E, 0x0B)   # orange rect (process)
ORANGE_DEEP = RGBColor(0xE2, 0x6D, 0x0F)   # diamond-ish accent
NOTE_FILL   = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT   = RGBColor(0x1A, 0x20, 0x2C)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY        = RGBColor(0x4A, 0x55, 0x68)
GRAY_LIGHT  = RGBColor(0xE2, 0xE8, 0xF0)

# Department palette (matches platform role theme)
DEPT_GREEN  = RGBColor(0x15, 0x80, 0x3D)
DEPT_RED    = RGBColor(0xB9, 0x1C, 0x1C)
DEPT_AMBER  = RGBColor(0xA1, 0x62, 0x07)
DEPT_BLUE   = RGBColor(0x1E, 0x40, 0xAF)
DEPT_GREY   = RGBColor(0x47, 0x55, 0x69)

# ─────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────
def add_shape(shape_type, left, top, width, height, *, fill=None,
              text="", text_color=WHITE, font_size=12, bold=True,
              line_color=None, line_width=1.5):
    sh = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill is not None:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line_color is not None:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(line_width)
    else:
        sh.line.color.rgb = DARK_TEXT
        sh.line.width = Pt(line_width)

    tf = sh.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Text
    tf.text = ""
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return sh

def add_multiline_shape(shape_type, left, top, width, height, *, fill,
                        title, bullets, title_color=WHITE,
                        body_color=DARK_TEXT, title_size=14, body_size=10):
    """Box with bold title on top + small bullet list below."""
    sh = slide.shapes.add_shape(shape_type, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = DARK_TEXT
    sh.line.width = Pt(1.25)

    tf = sh.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.text = ""

    # Title paragraph
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.name = "Calibri"
    r.font.size = Pt(title_size)
    r.font.bold = True
    r.font.color.rgb = title_color

    # Bullet paragraphs
    for line in bullets:
        bp = tf.add_paragraph()
        bp.alignment = PP_ALIGN.LEFT
        br = bp.add_run()
        br.text = "•  " + line
        br.font.name = "Calibri"
        br.font.size = Pt(body_size)
        br.font.bold = False
        br.font.color.rgb = body_color
    return sh

def add_arrow(start_x, start_y, end_x, end_y, *, color=DARK_TEXT, width=1.75):
    """Straight-line arrow connector from (start_x, start_y) to (end_x, end_y)."""
    line = slide.shapes.add_connector(2,   # 2 = STRAIGHT (with arrows we'll set below)
                                       Inches(start_x), Inches(start_y),
                                       Inches(end_x),   Inches(end_y))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    # Set end-arrow head via XML
    from pptx.oxml.ns import qn
    from lxml import etree
    spPr = line.line._get_or_add_ln()
    # Clear existing tail/head arrows and add a triangle head at the end
    for child in list(spPr):
        if child.tag in (qn("a:tailEnd"), qn("a:headEnd")):
            spPr.remove(child)
    tail = etree.SubElement(spPr, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w",   "med")
    tail.set("h",   "med")
    return line

def add_label(left, top, width, height, text, *, color=DARK_TEXT,
              size=11, bold=False, italic=False, align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

# ─────────────────────────────────────────────────────────────────────
# TITLE BAR
# ─────────────────────────────────────────────────────────────────────
title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0), Inches(0),
                                    prs.slide_width, Inches(0.9))
title_bar.fill.solid()
title_bar.fill.fore_color.rgb = RGBColor(0x4A, 0xA8, 0xD8)
title_bar.line.fill.background()
tf = title_bar.text_frame
tf.text = ""
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "MES Data Flow"
r.font.name = "Calibri"
r.font.size = Pt(36)
r.font.bold = True
r.font.color.rgb = WHITE

# ─────────────────────────────────────────────────────────────────────
# LAYER 1 — EOL MACHINE  (green oval, top centre)
# ─────────────────────────────────────────────────────────────────────
EOL_W, EOL_H = 3.4, 1.0
EOL_LEFT = (16.5 - EOL_W) / 2     # centred
EOL_TOP  = 1.3
add_shape(MSO_SHAPE.OVAL, Inches(EOL_LEFT), Inches(EOL_TOP),
          Inches(EOL_W), Inches(EOL_H),
          fill=GREEN_LIGHT, text="EOL Machine",
          text_color=WHITE, font_size=20, bold=True, line_width=2)

# Subtitle below the oval
add_label(EOL_LEFT - 0.8, EOL_TOP + EOL_H + 0.05, EOL_W + 1.6, 0.3,
          "(Main Assembly + 3 sub-stations)",
          color=GRAY, size=10, italic=True)

# ─────────────────────────────────────────────────────────────────────
# LAYER 2 — PLC / CAMERA / ANDON BOX  (rounded rectangles)
# ─────────────────────────────────────────────────────────────────────
ATTACH_TOP = 3.0
ATTACH_H   = 2.7
ATTACH_W   = 3.6
GAP        = 0.6
ATTACH_LEFTS = [
    (16.5 - 3 * ATTACH_W - 2 * GAP) / 2,
]
ATTACH_LEFTS.append(ATTACH_LEFTS[0] + ATTACH_W + GAP)
ATTACH_LEFTS.append(ATTACH_LEFTS[1] + ATTACH_W + GAP)

attach_specs = [
    ("PLC", BLUE_DARK, [
        "status",
        "OK / NG counts",
        "cycle time",
        "model running",
        "Poka-Yoke checks",
        "per-sub-station counts",
    ]),
    ("Camera", ORANGE, [
        "live video stream",
        "per-cycle recording",
    ]),
    ("Andon Box", ORANGE_DEEP, [
        "visual alert state",
        "operator events",
        "audible alarm",
    ]),
]

attach_centres = []
for left, (title, color, bullets) in zip(ATTACH_LEFTS, attach_specs):
    add_multiline_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(left), Inches(ATTACH_TOP),
                        Inches(ATTACH_W), Inches(ATTACH_H),
                        fill=color, title=title, bullets=bullets,
                        title_size=16, body_size=11)
    attach_centres.append(left + ATTACH_W/2)

# Arrows from EOL Machine to each attached box
eol_bottom_x = EOL_LEFT + EOL_W/2
eol_bottom_y = EOL_TOP + EOL_H
for cx in attach_centres:
    add_arrow(eol_bottom_x, eol_bottom_y + 0.05, cx, ATTACH_TOP)

# ─────────────────────────────────────────────────────────────────────
# LAYER 3 — COLLECTOR  (blue rectangle, full-width-ish)
# ─────────────────────────────────────────────────────────────────────
COLL_TOP = 6.1
COLL_H   = 1.0
COLL_W   = 6.5
COLL_LEFT = (16.5 - COLL_W) / 2
add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
          Inches(COLL_LEFT), Inches(COLL_TOP),
          Inches(COLL_W), Inches(COLL_H),
          fill=BLUE_DARK, text="Collector",
          text_color=WHITE, font_size=22, bold=True, line_width=2)

add_label(COLL_LEFT - 0.5, COLL_TOP + COLL_H + 0.02, COLL_W + 1.0, 0.3,
          "reads every machine continuously   ·   logs all events to the database",
          color=GRAY, size=10, italic=True)

# Arrows: each attached box → collector top
coll_top_x = COLL_LEFT + COLL_W/2
for cx, left in zip(attach_centres, ATTACH_LEFTS):
    add_arrow(cx, ATTACH_TOP + ATTACH_H, coll_top_x, COLL_TOP)

# ─────────────────────────────────────────────────────────────────────
# LAYER 4 — STREAMS + SIDE PANEL
# ─────────────────────────────────────────────────────────────────────
STREAM_TOP = 7.9
STREAM_H   = 1.5
STREAM_W   = 2.1
STREAM_GAP = 0.2

stream_specs = [
    ("Live OEE",        ORANGE,      ["plan vs actual", "hourly rollup", "loss breakdown"]),
    ("Breakdown",       ORANGE,      ["slip + cause", "duration", "MTBF / MTTR"]),
    ("Poka-Yoke",       ORANGE,      ["bypass + stuck", "email escalation"]),
    ("Process",         ORANGE,      ["per-min samples", "vs target", "→ side panel"]),
    ("Deviation",       ORANGE,      ["Maint → Quality", "approval flow"]),
]

# 5 streams + 1 side panel laid out across the slide
total_streams_w = 5 * STREAM_W + 4 * STREAM_GAP
SIDE_W = 2.8
SIDE_GAP = 0.5
content_w = total_streams_w + SIDE_GAP + SIDE_W
content_left = (16.5 - content_w) / 2

stream_centres = []
left_cursor = content_left
for title, color, bullets in stream_specs:
    add_multiline_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(left_cursor), Inches(STREAM_TOP),
                        Inches(STREAM_W), Inches(STREAM_H),
                        fill=color, title=title, bullets=bullets,
                        title_size=12, body_size=9)
    stream_centres.append(left_cursor + STREAM_W/2)
    left_cursor += STREAM_W + STREAM_GAP

# Side panel — taller, on the right
SIDE_LEFT = left_cursor - STREAM_GAP + SIDE_GAP
SIDE_TOP  = STREAM_TOP - 0.05
SIDE_H    = STREAM_H + 0.5
add_multiline_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(SIDE_LEFT), Inches(SIDE_TOP),
                    Inches(SIDE_W), Inches(SIDE_H),
                    fill=GRAY, title="Side Panel",
                    bullets=[
                        "Main Assembly",
                        "Upper Rail Greasing",
                        "Lock Bar Insert",
                        "Lower Rail Greasing",
                        "(per-min bar chart vs target)",
                    ],
                    title_size=12, body_size=9)

# Arrow from Process Tracking → Side panel
proc_idx = 3
proc_right_x = stream_centres[proc_idx] + STREAM_W/2
proc_mid_y   = STREAM_TOP + STREAM_H/2
add_arrow(proc_right_x, proc_mid_y, SIDE_LEFT, proc_mid_y, color=GRAY, width=1.25)

# Arrows: collector bottom → each stream
coll_bottom_x = COLL_LEFT + COLL_W/2
coll_bottom_y = COLL_TOP + COLL_H + 0.35   # leave space for caption
for cx in stream_centres:
    add_arrow(coll_bottom_x, coll_bottom_y, cx, STREAM_TOP)

# ─────────────────────────────────────────────────────────────────────
# LAYER 5 — DEPARTMENTS  (5 coloured ovals)
# ─────────────────────────────────────────────────────────────────────
DEPT_TOP = 10.0
DEPT_H   = 1.1
DEPT_W   = 2.4
DEPT_GAP = 0.45

dept_specs = [
    ("Production Team",  DEPT_GREEN),
    ("Maintenance Team", DEPT_RED),
    ("Quality Team",     DEPT_AMBER),
    ("Admin / Plant Head", DEPT_BLUE),
    ("Operator",         DEPT_GREY),
]
dept_total_w = 5 * DEPT_W + 4 * DEPT_GAP
dept_left = (16.5 - dept_total_w) / 2
dept_centres = []
left = dept_left
for label, color in dept_specs:
    add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(DEPT_TOP),
              Inches(DEPT_W), Inches(DEPT_H),
              fill=color, text=label,
              text_color=WHITE, font_size=12.5, bold=True, line_width=1.75)
    dept_centres.append(left + DEPT_W/2)
    left += DEPT_W + DEPT_GAP

# Arrows from each stream to its target dept
# Mapping: stream index → list of dept indices
stream_to_dept = {
    0: [0, 3],          # Live OEE → Production, Admin
    1: [0, 1, 3],       # Breakdown → Production banner, Maintenance, Admin
    2: [1, 2, 3],       # Poka-Yoke → Maintenance, Quality, Admin
    3: [0, 3],          # Process → Production, Admin
    4: [1, 2, 3],       # Deviation → Maintenance raise, Quality approve, Admin
}
for s_idx, dept_idxs in stream_to_dept.items():
    sx = stream_centres[s_idx]
    sy = STREAM_TOP + STREAM_H
    for d_idx in dept_idxs:
        dx = dept_centres[d_idx]
        dy = DEPT_TOP
        add_arrow(sx, sy, dx, dy, color=GRAY_LIGHT, width=0.75)

# Dashed line from Production-related streams to Operator (restricted)
add_arrow(stream_centres[0], STREAM_TOP + STREAM_H,
           dept_centres[4], DEPT_TOP, color=GRAY, width=0.75)

# ─────────────────────────────────────────────────────────────────────
# Footer
# ─────────────────────────────────────────────────────────────────────
add_label(0, 11.25, 16.5, 0.3,
          "Solid arrows = primary data flow.   Operators see only their own assigned lines.   Internal — Toyota Boshoku Device India, Bawal.",
          color=GRAY, size=9, italic=True)

# ─────────────────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────────────────
import time
out_path = r"D:\EOL\EOL\Deep (2)\Deep\MES_Data_Flow.pptx"
try:
    prs.save(out_path)
except PermissionError:
    out_path = out_path.replace(".pptx", f"_{int(time.time())}.pptx")
    prs.save(out_path)
    print(f"WARNING: original was locked — saved as: {out_path}")
print(f"PowerPoint written: {out_path}")
